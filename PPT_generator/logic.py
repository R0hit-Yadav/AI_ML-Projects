import pptx
from pptx.util import Pt
import ollama
import os
import textwrap
import time

# Font size constants
TITLE_FONT_SIZE = Pt(32)
CONTENT_FONT_SIZE = Pt(18)
FONT_NAME = "Calibri"

def generate_text(prompt):
    # Take Response From Llama
    response = ollama.chat(model="llama3.2", messages=[ 
        {"role": "system", "content": "You are a helpful assistant."},
        {"role": "user", "content": prompt}
    ])
    
    if 'message' in response and 'content' in response['message']:
        return response['message']['content'].strip()
    return 'No text found'

#genrate numebr of slide and slide title for given topic
def generate_slide_titles(topic,n_slide):
    prompt = f"Generate {n_slide} concise slide titles for the topic '{topic}'. Each title should be short and engaging. and remove ** from start and end of title and add emoji and sign for better formating"
    response = generate_text(prompt)
    return [title.strip() for title in response.split("\n") if title.strip()][:n_slide]

#genrate each slide contetnt using title of every slide 
def generate_slide_content(slide_title):
    prompt = f"Generate a well-structured slide content for '{slide_title}'.Provide 3-4 concise bullet points in proper formating. and remove 'Here is a well-structured slide content' line and remove extra ** from contant and provide bullet point and remove ** from bullet points and add emoji and sign for better formating  "
    response = generate_text(prompt)

    # Extracting bullet points properly
    lines = response.split("\n")
    bullet_points = [line.strip("-•").strip() for line in lines if line.strip() and not line.lower().startswith("title")]
    
    # define number of bullets
    return bullet_points[:4]

#create ppt
def create_presentation(topic, slide_titles, slide_contents):
    prs = pptx.Presentation()

    # formating of title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_text_frame = title_slide.shapes.title.text_frame
    title_text_frame.paragraphs[0].font.size = TITLE_FONT_SIZE
    title_text_frame.paragraphs[0].font.name = FONT_NAME


    for slide_title, slide_content in zip(slide_titles, slide_contents):
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_title

        # Format title
        title_tf = slide.shapes.title.text_frame
        title_tf.paragraphs[0].font.size = TITLE_FONT_SIZE
        title_tf.paragraphs[0].font.name = FONT_NAME
        title_tf.paragraphs[0].bold = True

        content_shape = slide.shapes.placeholders[1]
        content_text_frame = content_shape.text_frame
        content_text_frame.clear()

        for point in slide_content:
            para = content_text_frame.add_paragraph()
            para.text = textwrap.fill(point, width=80)  # Wrap text for better readability
            para.font.size = CONTENT_FONT_SIZE
            para.font.name = FONT_NAME
            para.space_after = Pt(10)  # Adjust spacing

    #stored every ppt on folder 
    if not os.path.exists("generated_ppt"):
        os.makedirs("generated_ppt")

    ppt_path = f"generated_ppt/{topic}_presentation.pptx"
    prs.save(ppt_path)
    return ppt_path

